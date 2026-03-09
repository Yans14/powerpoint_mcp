from __future__ import annotations

from typing import Any

from lxml import etree
from pptx import Presentation


def _extract_theme_colors(prs: Presentation) -> dict[str, str]:
    colors: dict[str, str] = {}
    try:
        master = prs.slide_masters[0]
        for rel in master.part.rels.values():
            if "theme" in rel.reltype:
                theme_part = rel.target_part
                root = etree.fromstring(theme_part.blob)
                for clr_elem in root.findall(".//{*}srgbClr"):
                    val = clr_elem.get("val", "")
                    if val:
                        parent_tag = clr_elem.getparent().tag.split("}")[-1]
                        colors[parent_tag] = val.upper()
                break
    except Exception:
        pass
    return colors


class TemplateConformanceChecker:
    def check(
        self,
        prs: Presentation,
        template_path: str,
        check_theme: bool = True,
        check_fonts: bool = True,
        check_layouts: bool = True,
    ) -> dict[str, Any]:
        template_prs = Presentation(template_path)
        issues: list[dict[str, Any]] = []
        score_components: list[float] = []

        if check_theme:
            src_colors = _extract_theme_colors(prs)
            tpl_colors = _extract_theme_colors(template_prs)
            if src_colors and tpl_colors:
                matching = sum(1 for k, v in tpl_colors.items() if src_colors.get(k) == v)
                theme_score = matching / len(tpl_colors) if tpl_colors else 1.0
                score_components.append(theme_score)
                if theme_score < 0.8:
                    issues.append(
                        {
                            "issue_type": "theme_color_mismatch",
                            "description": f"Theme colors match {matching}/{len(tpl_colors)} template colors.",
                            "severity": "warning",
                            "source_colors": src_colors,
                            "template_colors": tpl_colors,
                        }
                    )

        if check_layouts:
            tpl_layout_names = {layout.name for layout in template_prs.slide_layouts}
            used_layouts = {prs.slides[i].slide_layout.name for i in range(len(prs.slides))}
            missing = used_layouts - tpl_layout_names
            layout_score = 1.0 - (len(missing) / max(len(used_layouts), 1))
            score_components.append(layout_score)
            for layout_name in missing:
                issues.append(
                    {
                        "issue_type": "missing_layout",
                        "description": f"Layout '{layout_name}' (used in presentation) does not exist in template.",
                        "severity": "critical",
                    }
                )

        conformance_score = sum(score_components) / len(score_components) if score_components else 0.5

        return {
            "conformance_score": round(conformance_score, 3),
            "issues": issues,
            "summary": {
                "total_issues": len(issues),
                "critical": sum(1 for i in issues if i.get("severity") == "critical"),
                "warnings": sum(1 for i in issues if i.get("severity") == "warning"),
            },
        }
