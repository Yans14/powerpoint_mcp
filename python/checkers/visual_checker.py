from __future__ import annotations

from collections import Counter
from typing import Any

from pptx import Presentation


def _get_run_font(run) -> str | None:
    try:
        return run.font.name
    except Exception:
        return None


def _get_run_color(run) -> str | None:
    try:
        color = run.font.color
        if color and color.type is not None and color.rgb is not None:
            return str(color.rgb).upper()
    except Exception:
        pass
    return None


class VisualConsistencyChecker:
    def check(
        self,
        prs: Presentation,
        slide_indices: list[int],
        check_fonts: bool = True,
        check_colors: bool = True,
        check_sizes: bool = True,
    ) -> dict[str, Any]:
        all_fonts: list[str] = []
        all_colors: list[str] = []
        all_sizes: list[float] = []
        per_slide_fonts: dict[int, set[str]] = {}
        per_slide_colors: dict[int, set[str]] = {}
        issues: list[dict[str, Any]] = []

        for idx in slide_indices:
            slide = prs.slides[idx - 1]
            slide_fonts: set[str] = set()
            slide_colors: set[str] = set()

            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if check_fonts:
                            font = _get_run_font(run)
                            if font:
                                all_fonts.append(font)
                                slide_fonts.add(font)
                        if check_colors:
                            color = _get_run_color(run)
                            if color:
                                all_colors.append(color)
                                slide_colors.add(color)
                        if check_sizes:
                            try:
                                if run.font.size:
                                    all_sizes.append(run.font.size / 12700)
                            except Exception:
                                pass

            per_slide_fonts[idx] = slide_fonts
            per_slide_colors[idx] = slide_colors

        font_counts = Counter(all_fonts)
        color_counts = Counter(all_colors)

        dominant_font = font_counts.most_common(1)[0][0] if font_counts else None

        if check_fonts and dominant_font:
            for idx, fonts in per_slide_fonts.items():
                non_dominant = fonts - {dominant_font}
                if non_dominant:
                    issues.append(
                        {
                            "slide_index": idx,
                            "issue_type": "inconsistent_font",
                            "description": f"Slide {idx} uses fonts {non_dominant} (dominant: {dominant_font})",
                            "severity": "warning",
                        }
                    )

        return {
            "font_report": {
                "fonts_used": dict(font_counts),
                "dominant_font": dominant_font,
                "slide_font_map": {k: list(v) for k, v in per_slide_fonts.items()},
            },
            "color_report": {
                "colors_used": dict(color_counts),
                "slide_color_map": {k: list(v) for k, v in per_slide_colors.items()},
            },
            "issues": issues,
            "summary": {
                "total_issues": len(issues),
                "unique_fonts": len(font_counts),
                "unique_colors": len(color_counts),
                "slides_checked": len(slide_indices),
            },
        }
