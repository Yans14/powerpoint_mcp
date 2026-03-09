from __future__ import annotations

from typing import Any

from pptx import Presentation


def _slide_text_map(slide) -> dict[str, str]:
    result: dict[str, str] = {}
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            result[shape.name] = shape.text_frame.text.strip()
    return result


class PresentationDiffer:
    def diff(
        self,
        prs_a: Presentation,
        prs_b: Presentation,
        deep_diff: bool = True,
    ) -> dict[str, Any]:
        count_a = len(prs_a.slides)
        count_b = len(prs_b.slides)

        added_slides: list[int] = []
        removed_slides: list[int] = []
        modified_slides: list[dict[str, Any]] = []

        for i in range(min(count_a, count_b)):
            slide_a = prs_a.slides[i]
            slide_b = prs_b.slides[i]
            changes: list[dict[str, Any]] = []

            if deep_diff:
                text_a = _slide_text_map(slide_a)
                text_b = _slide_text_map(slide_b)
                all_shapes = set(text_a.keys()) | set(text_b.keys())
                for shape_name in all_shapes:
                    val_a = text_a.get(shape_name)
                    val_b = text_b.get(shape_name)
                    if val_a != val_b:
                        changes.append(
                            {
                                "shape_name": shape_name,
                                "before": val_a,
                                "after": val_b,
                            }
                        )

            if slide_a.slide_layout.name != slide_b.slide_layout.name:
                changes.append(
                    {
                        "field": "layout",
                        "before": slide_a.slide_layout.name,
                        "after": slide_b.slide_layout.name,
                    }
                )

            if changes:
                modified_slides.append({"slide_index": i + 1, "changes": changes})

        for i in range(count_b - count_a):
            added_slides.append(count_a + i + 1)

        for i in range(count_a - count_b):
            removed_slides.append(count_b + i + 1)

        return {
            "slide_count_diff": {"before": count_a, "after": count_b},
            "added_slides": added_slides,
            "removed_slides": removed_slides,
            "modified_slides": modified_slides,
            "summary": (
                f"{count_b - count_a:+d} slides, "
                f"{len(modified_slides)} modified slides, "
                f"{sum(len(s['changes']) for s in modified_slides)} field changes."
            ),
        }
