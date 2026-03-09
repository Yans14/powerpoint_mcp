from __future__ import annotations

import re
from typing import Any

from pptx import Presentation

_DEFAULT_TEXT_PATTERNS = [
    re.compile(r"click to (add|edit|type)", re.IGNORECASE),
    re.compile(r"click here to (add|edit|type)", re.IGNORECASE),
    re.compile(r"^(title\s*\d*|subtitle|text box|content placeholder)$", re.IGNORECASE),
]


class ContentChecker:
    def check(
        self,
        prs: Presentation,
        check_empty: bool = True,
        check_default_text: bool = True,
        slide_indices: list[int] | None = None,
    ) -> dict[str, Any]:
        indices = slide_indices or list(range(1, len(prs.slides) + 1))
        empty_placeholders: list[dict[str, Any]] = []
        default_text_found: list[dict[str, Any]] = []
        issues: list[dict[str, Any]] = []

        for idx in indices:
            slide = prs.slides[idx - 1]
            for shape in slide.placeholders:
                ph_name = shape.name
                ph_type = str(shape.placeholder_format.type) if shape.placeholder_format else ""

                if check_empty:
                    is_empty = True
                    if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                        is_empty = False
                    if hasattr(shape, "has_picture") and shape.shape_type == 13:
                        is_empty = False
                    if is_empty:
                        entry = {
                            "slide_index": idx,
                            "placeholder_name": ph_name,
                            "placeholder_type": ph_type,
                        }
                        empty_placeholders.append(entry)
                        issues.append(
                            {
                                "slide_index": idx,
                                "issue_type": "empty_placeholder",
                                "description": f"Placeholder '{ph_name}' on slide {idx} is empty.",
                                "severity": "warning",
                            }
                        )

                if check_default_text and getattr(shape, "has_text_frame", False):
                    text = shape.text_frame.text.strip()
                    for pattern in _DEFAULT_TEXT_PATTERNS:
                        if pattern.search(text):
                            default_text_found.append(
                                {
                                    "slide_index": idx,
                                    "placeholder_name": ph_name,
                                    "text": text,
                                }
                            )
                            issues.append(
                                {
                                    "slide_index": idx,
                                    "issue_type": "default_text",
                                    "description": f"Placeholder '{ph_name}' on slide {idx} contains default text: '{text[:50]}'",
                                    "severity": "warning",
                                }
                            )
                            break

        return {
            "empty_placeholders": empty_placeholders,
            "default_text_remaining": default_text_found,
            "issues": issues,
            "summary": {
                "total_issues": len(issues),
                "empty_count": len(empty_placeholders),
                "default_text_count": len(default_text_found),
                "slides_checked": len(indices),
            },
        }
