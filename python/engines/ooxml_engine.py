from __future__ import annotations

import base64
import copy
import os
import re
import shutil
import subprocess
import sys
import tempfile
import uuid
import zipfile
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from engines.base import BaseEngine, EngineSession
from errors import BridgeError, ensure
from utils.colors import normalize_color
from utils.paths import validate_existing_file, validate_output_file
from utils.units import emu_to_inches, to_emu

_COLOR_NAMESPACES = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}


class OOXMLEngine(BaseEngine):
    name = "OOXML"

    def __init__(self, metadata: dict[str, Any]) -> None:
        self.metadata = metadata
        self.sessions: dict[str, EngineSession] = {}

    def get_engine_info(self, params: dict[str, Any] | None = None) -> dict[str, Any]:
        return {
            "engine": self.name,
            **self.metadata,
        }

    def create_presentation(self, params: dict[str, Any]) -> dict[str, Any]:
        width = params.get("width")
        height = params.get("height")
        template_path = params.get("template_path")

        if template_path:
            template = validate_existing_file(str(template_path), expected_suffixes=(".pptx", ".potx"))
            fd, tmp_path = tempfile.mkstemp(suffix=".pptx", prefix="pptx-session-")
            os.close(fd)
            Path(tmp_path).unlink(missing_ok=True)
            shutil.copy2(template, tmp_path)
            prs = Presentation(str(tmp_path))
            original_path = str(template)
        else:
            prs = Presentation()
            if width is not None:
                prs.slide_width = to_emu(width)
            if height is not None:
                prs.slide_height = to_emu(height)
            fd, tmp_path = tempfile.mkstemp(suffix=".pptx", prefix="pptx-session-")
            os.close(fd)
            Path(tmp_path).unlink(missing_ok=True)
            prs.save(tmp_path)
            prs = Presentation(tmp_path)
            original_path = ""

        session_id = str(uuid.uuid4())
        session = EngineSession(
            id=session_id,
            original_path=original_path,
            working_path=tmp_path,
            engine=self.name,
            dirty=True,
            extra={"prs": prs},
        )
        self.sessions[session_id] = session

        return {
            "success": True,
            "presentation_id": session_id,
            "engine": self.name,
            "slide_width_emu": prs.slide_width,
            "slide_height_emu": prs.slide_height,
            "slide_count": len(prs.slides),
            "presentation_state": self.get_presentation_state({"presentation_id": session_id}),
        }

    def open_presentation(self, params: dict[str, Any]) -> dict[str, Any]:
        source = validate_existing_file(str(params["file_path"]), expected_suffixes=(".pptx", ".potx"))

        fd, tmp_path = tempfile.mkstemp(suffix=".pptx", prefix="pptx-session-")
        os.close(fd)
        Path(tmp_path).unlink(missing_ok=True)
        shutil.copy2(source, tmp_path)

        session_id = str(uuid.uuid4())
        prs = Presentation(tmp_path)
        session = EngineSession(
            id=session_id,
            original_path=str(source),
            working_path=tmp_path,
            engine=self.name,
            dirty=False,
            extra={"prs": prs},
        )
        self.sessions[session_id] = session

        return {
            "success": True,
            "presentation_id": session_id,
            "slide_count": len(prs.slides),
            "layout_names": [layout.name for layout in prs.slide_layouts],
            "presentation_state": self.get_presentation_state({"presentation_id": session_id}),
        }

    def save_presentation(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        output_path = validate_output_file(str(params["output_path"]), expected_suffixes=(".pptx",))
        prs = self._prs(session)
        prs.save(session.working_path)
        shutil.copy2(session.working_path, output_path)
        session.dirty = False

        return {
            "success": True,
            "presentation_id": session.id,
            "saved_path": str(output_path),
            "file_size_bytes": output_path.stat().st_size,
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def close_presentation(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        self.sessions.pop(session.id, None)
        try:
            Path(session.working_path).unlink(missing_ok=True)
        except Exception:
            pass

        return {
            "success": True,
            "presentation_id": session.id,
            "closed": True,
        }

    def list_open_presentations(self, params: dict[str, Any] | None = None) -> dict[str, Any]:
        return {
            "presentations": [
                {
                    "presentation_id": session.id,
                    "original_path": session.original_path,
                    "working_path": session.working_path,
                    "engine": session.engine,
                    "dirty": session.dirty,
                    "opened_at": session.opened_at.isoformat(),
                }
                for session in self.sessions.values()
            ]
        }

    def get_presentation_state(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        prs = self._prs(session)

        slides = []
        for index, slide in enumerate(prs.slides, start=1):
            slides.append(
                {
                    "index": index,
                    "title": self._slide_title(slide),
                    "layout": slide.slide_layout.name,
                    "shape_count": len(slide.shapes),
                }
            )

        return {
            "presentation_id": session.id,
            "engine": self.name,
            "slide_count": len(prs.slides),
            "slides": slides,
        }

    def get_layouts(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        prs = self._prs(session)

        layouts: list[dict[str, Any]] = []
        for idx, layout in enumerate(prs.slide_layouts, start=1):
            placeholder_types = []
            for placeholder in layout.placeholders:
                placeholder_types.append(str(placeholder.placeholder_format.type))

            layouts.append(
                {
                    "index": idx,
                    "name": layout.name,
                    "placeholder_types": placeholder_types,
                }
            )

        return {
            "presentation_id": session.id,
            "layouts": layouts,
        }

    def get_layout_detail(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        prs = self._prs(session)

        layout_name = str(params["layout_name"])
        layout = self._find_layout(prs, layout_name)

        placeholders = [self._placeholder_payload(placeholder) for placeholder in layout.placeholders]

        return {
            "presentation_id": session.id,
            "layout_name": layout.name,
            "placeholder_count": len(placeholders),
            "placeholders": placeholders,
        }

    def get_masters(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        prs = self._prs(session)

        masters = []
        for idx, master in enumerate(prs.slide_masters, start=1):
            masters.append(
                {
                    "index": idx,
                    "name": getattr(master, "name", f"Master {idx}"),
                    "layout_count": len(master.slide_layouts),
                }
            )

        return {
            "presentation_id": session.id,
            "masters": masters,
        }

    def get_theme(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        theme_data = self._extract_theme(Path(session.working_path))

        return {
            "presentation_id": session.id,
            "theme": theme_data,
        }

    def get_slide(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text

        return {
            "presentation_id": session.id,
            "slide_index": slide_index,
            "title": self._slide_title(slide),
            "layout": slide.slide_layout.name,
            "shapes": [self._shape_payload(shape) for shape in slide.shapes],
            "placeholders": [self._placeholder_payload(placeholder) for placeholder in slide.placeholders],
            "notes": notes_text,
        }

    def add_slide(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        prs = self._prs(session)
        layout_name = str(params["layout_name"])
        layout = self._find_layout(prs, layout_name)

        new_slide = prs.slides.add_slide(layout)
        new_index = len(prs.slides)

        requested_position = params.get("position")
        if requested_position is not None:
            position = int(requested_position)
            ensure(position >= 1 and position <= len(prs.slides), "validation_error", "position is out of bounds")
            if position != new_index:
                order = list(range(1, len(prs.slides) + 1))
                moved = order.pop(new_index - 1)
                order.insert(position - 1, moved)
                self._reorder_slide_ids(prs, order)
                new_index = position

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "added_slide_index": new_index,
            "layout_name": layout_name,
            "placeholders": [self._placeholder_payload(placeholder) for placeholder in new_slide.placeholders],
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def duplicate_slide(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        prs = self._prs(session)
        source_index = int(params["source_index"])
        source_slide = self._get_slide_obj(session, source_index)

        duplicate = prs.slides.add_slide(source_slide.slide_layout)

        # Copy non-placeholder shapes directly.
        for shape in source_slide.shapes:
            if getattr(shape, "is_placeholder", False):
                continue
            new_el = copy.deepcopy(shape.element)
            duplicate.shapes._spTree.insert_element_before(new_el, "p:extLst")

        # Copy placeholder text where names match.
        for src_placeholder in source_slide.placeholders:
            try:
                dst_placeholder = self._find_placeholder(duplicate, src_placeholder.name)
            except BridgeError:
                continue
            if (
                hasattr(src_placeholder, "has_text_frame")
                and src_placeholder.has_text_frame
                and dst_placeholder.has_text_frame
            ):
                dst_placeholder.text = src_placeholder.text

        if source_slide.has_notes_slide and source_slide.notes_slide.notes_text_frame:
            duplicate.notes_slide.notes_text_frame.text = source_slide.notes_slide.notes_text_frame.text

        duplicate_index = len(prs.slides)
        requested_position = params.get("target_position")
        if requested_position is not None:
            position = int(requested_position)
            ensure(
                position >= 1 and position <= len(prs.slides), "validation_error", "target_position is out of bounds"
            )
            if position != duplicate_index:
                order = list(range(1, len(prs.slides) + 1))
                moved = order.pop(duplicate_index - 1)
                order.insert(position - 1, moved)
                self._reorder_slide_ids(prs, order)
                duplicate_index = position

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "source_index": source_index,
            "duplicated_slide_index": duplicate_index,
            "warning": "OOXML duplicate copies placeholders and non-placeholder shapes; complex animations may require COM mode.",
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def delete_slide(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        prs = self._prs(session)
        slide_index = int(params["slide_index"])
        self._assert_slide_index(prs, slide_index)

        sld_id_lst = prs.slides._sldIdLst
        slide_ids = list(sld_id_lst)
        slide_id = slide_ids[slide_index - 1]
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        sld_id_lst.remove(slide_id)

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "deleted_slide_index": slide_index,
            "warning": f"Slide indices above {slide_index} shifted down by 1.",
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def reorder_slides(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        prs = self._prs(session)

        new_order = [int(value) for value in params["new_order"]]
        self._validate_order(prs, new_order)
        self._reorder_slide_ids(prs, new_order)

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "new_order": new_order,
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def move_slide(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        prs = self._prs(session)

        from_index = int(params["from_index"])
        to_index = int(params["to_index"])
        self._assert_slide_index(prs, from_index)
        self._assert_slide_index(prs, to_index)

        order = list(range(1, len(prs.slides) + 1))
        moved = order.pop(from_index - 1)
        order.insert(to_index - 1, moved)
        self._reorder_slide_ids(prs, order)

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "from_index": from_index,
            "to_index": to_index,
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def set_slide_background(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        prs = self._prs(session)
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        warnings: list[str] = []

        color_hex = params.get("color_hex")
        image_path = params.get("image_path")
        grad_start = params.get("gradient_start_color_hex")
        grad_end = params.get("gradient_end_color_hex")

        ensure(
            any([color_hex, image_path, grad_start and grad_end]), "validation_error", "No background input provided."
        )

        if color_hex:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(normalize_color(str(color_hex)))

        if image_path:
            image = validate_existing_file(str(image_path), expected_suffixes=(".png", ".jpg", ".jpeg", ".bmp", ".gif"))
            picture = slide.shapes.add_picture(str(image), 0, 0, prs.slide_width, prs.slide_height)
            self._send_shape_to_back(slide, picture)
            warnings.append(
                "OOXML fallback inserts a full-slide picture shape at back. Use COM mode for native background picture fidelity."
            )

        if grad_start and grad_end:
            fill = slide.background.fill
            try:
                fill.gradient()
                grad_stops = fill.gradient_stops
                grad_stops[0].color.rgb = RGBColor.from_string(normalize_color(str(grad_start)))
                grad_stops[-1].color.rgb = RGBColor.from_string(normalize_color(str(grad_end)))
            except Exception:
                warnings.append("Gradient background is partially supported in OOXML mode.")

        self._persist(session)

        payload: dict[str, Any] = {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }
        if warnings:
            payload["warning"] = " ".join(warnings)
        return payload

    def get_slide_snapshot(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        width_px = int(params.get("width_px") or 1280)
        self._assert_slide_index(self._prs(session), slide_index)

        soffice = shutil.which("soffice")
        pdftoppm = shutil.which("pdftoppm")

        if not soffice or not pdftoppm:
            raise BridgeError(
                code="dependency_missing",
                message="pptx_get_slide_snapshot requires LibreOffice (soffice) and poppler (pdftoppm) in OOXML mode.",
                details={
                    "missing": [
                        dependency
                        for dependency, ok in {
                            "soffice": bool(soffice),
                            "pdftoppm": bool(pdftoppm),
                        }.items()
                        if not ok
                    ],
                    "install_hint": "Install LibreOffice and poppler-utils to enable slide snapshots.",
                },
            )

        self._persist(session)

        with tempfile.TemporaryDirectory(prefix="pptx-snapshot-") as tmp_dir:
            source_path = Path(session.working_path)
            pdf_out_dir = Path(tmp_dir)

            convert_result = subprocess.run(
                [
                    str(soffice),
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(pdf_out_dir),
                    str(source_path),
                ],
                capture_output=True,
                text=True,
                check=False,
            )

            if convert_result.returncode != 0:
                raise BridgeError(
                    code="engine_error",
                    message="LibreOffice PDF export failed.",
                    details={
                        "stdout": convert_result.stdout,
                        "stderr": convert_result.stderr,
                    },
                )

            pdf_path = pdf_out_dir / f"{source_path.stem}.pdf"
            ensure(pdf_path.exists(), "engine_error", "Expected PDF output was not generated.")

            image_prefix = pdf_out_dir / "slide"
            render_result = subprocess.run(
                [
                    str(pdftoppm),
                    "-f",
                    str(slide_index),
                    "-l",
                    str(slide_index),
                    "-jpeg",
                    "-singlefile",
                    "-scale-to-x",
                    str(width_px),
                    str(pdf_path),
                    str(image_prefix),
                ],
                capture_output=True,
                text=True,
                check=False,
            )
            if render_result.returncode != 0:
                raise BridgeError(
                    code="engine_error",
                    message="pdftoppm rendering failed.",
                    details={
                        "stdout": render_result.stdout,
                        "stderr": render_result.stderr,
                    },
                )

            image_path = Path(f"{image_prefix}.jpg")
            ensure(image_path.exists(), "engine_error", "Slide snapshot image was not generated.")
            encoded = base64.b64encode(image_path.read_bytes()).decode("ascii")

        return {
            "presentation_id": session.id,
            "slide_index": slide_index,
            "mime_type": "image/jpeg",
            "snapshot_base64": encoded,
            "width_px": width_px,
        }

    def get_placeholders(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide = self._get_slide_obj(session, int(params["slide_index"]))

        placeholders = [self._placeholder_payload(placeholder) for placeholder in slide.placeholders]
        return {
            "presentation_id": session.id,
            "slide_index": int(params["slide_index"]),
            "placeholders": placeholders,
        }

    def set_placeholder_text(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)
        placeholder = self._find_placeholder(slide, str(params["placeholder_name"]))

        ensure(placeholder.has_text_frame, "conflict", "Target placeholder does not support text.")

        text_frame = placeholder.text_frame
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        paragraph.text = str(params.get("text_content", ""))

        alignment = params.get("alignment")
        if alignment and alignment in self._ALIGN_MAP:
            paragraph.alignment = self._ALIGN_MAP[alignment]

        if paragraph.runs:
            run = paragraph.runs[0]
            if params.get("font_name"):
                run.font.name = str(params["font_name"])
            if params.get("font_size_pt"):
                run.font.size = Pt(float(params["font_size_pt"]))
            if params.get("bold") is not None:
                run.font.bold = bool(params["bold"])
            if params.get("italic") is not None:
                run.font.italic = bool(params["italic"])
            if params.get("underline") is not None:
                run.font.underline = bool(params["underline"])
            if params.get("color_hex"):
                run.font.color.rgb = RGBColor.from_string(normalize_color(str(params["color_hex"])))

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "placeholder_name": str(params["placeholder_name"]),
            "text_content": str(params.get("text_content", "")),
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def set_placeholder_image(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)
        placeholder = self._find_placeholder(slide, str(params["placeholder_name"]))

        image_path = validate_existing_file(
            str(params["image_path"]),
            expected_suffixes=(".png", ".jpg", ".jpeg", ".bmp", ".gif"),
        )

        if hasattr(placeholder, "insert_picture"):
            placeholder.insert_picture(str(image_path))
        else:
            slide.shapes.add_picture(
                str(image_path), placeholder.left, placeholder.top, placeholder.width, placeholder.height
            )

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "placeholder_name": str(params["placeholder_name"]),
            "image_path": str(image_path),
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def clear_placeholder(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)
        placeholder_name = str(params["placeholder_name"])
        placeholder = self._find_placeholder(slide, placeholder_name)

        warning = ""
        if placeholder.has_text_frame:
            placeholder.text_frame.clear()
        else:
            warning = "Non-text placeholder reset is limited in OOXML mode."

        self._persist(session)

        payload: dict[str, Any] = {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "placeholder_name": placeholder_name,
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }
        if warning:
            payload["warning"] = warning
        return payload

    def get_placeholder_text(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide = self._get_slide_obj(session, int(params["slide_index"]))
        placeholder = self._find_placeholder(slide, str(params["placeholder_name"]))

        ensure(placeholder.has_text_frame, "conflict", "Target placeholder does not contain text.")

        paragraphs = []
        for paragraph in placeholder.text_frame.paragraphs:
            runs = []
            for run in paragraph.runs:
                color_hex = None
                try:
                    if run.font.color and run.font.color.rgb:
                        color_hex = str(run.font.color.rgb)
                except Exception:
                    color_hex = None

                runs.append(
                    {
                        "text": run.text,
                        "bold": bool(run.font.bold) if run.font.bold is not None else None,
                        "italic": bool(run.font.italic) if run.font.italic is not None else None,
                        "underline": bool(run.font.underline) if run.font.underline is not None else None,
                        "font_name": run.font.name,
                        "font_size_pt": float(run.font.size.pt) if run.font.size else None,
                        "color_hex": color_hex,
                    }
                )
            paragraphs.append({"text": paragraph.text, "runs": runs})

        return {
            "presentation_id": session.id,
            "slide_index": int(params["slide_index"]),
            "placeholder_name": str(params["placeholder_name"]),
            "paragraphs": paragraphs,
            "raw_text": placeholder.text,
        }

    # ------------------------------------------------------------------ #
    # Phase 3A — Rich text & content reading
    # ------------------------------------------------------------------ #

    _ALIGN_MAP = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }

    _REVERSE_ALIGN_MAP = {v: k for k, v in _ALIGN_MAP.items()}

    def set_placeholder_rich_text(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)
        placeholder = self._find_placeholder(slide, str(params["placeholder_name"]))

        ensure(placeholder.has_text_frame, "conflict", "Target placeholder does not support text.")

        paragraphs_data = params.get("paragraphs", [])
        ensure(len(paragraphs_data) > 0, "validation_error", "At least one paragraph is required.")

        self._write_paragraphs(placeholder.text_frame, paragraphs_data)

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "placeholder_name": str(params["placeholder_name"]),
            "paragraph_count": len(paragraphs_data),
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def add_text_box(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        left = to_emu(params["left"])
        top = to_emu(params["top"])
        width = to_emu(params["width"])
        height = to_emu(params["height"])

        from pptx.util import Emu as EmuVal

        text_box = slide.shapes.add_textbox(EmuVal(left), EmuVal(top), EmuVal(width), EmuVal(height))
        text_frame = text_box.text_frame

        # Support rich paragraphs or simple text_content
        paragraphs_data = params.get("paragraphs")
        if paragraphs_data:
            self._write_paragraphs(text_frame, paragraphs_data)
        else:
            text_content = str(params.get("text_content", ""))
            paragraph = text_frame.paragraphs[0]
            paragraph.text = text_content

            alignment = params.get("alignment")
            if alignment and alignment in self._ALIGN_MAP:
                paragraph.alignment = self._ALIGN_MAP[alignment]

            if paragraph.runs:
                run = paragraph.runs[0]
                self._apply_run_formatting(run, params)

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "shape_name": text_box.name,
            "shape_id": int(text_box.shape_id),
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def get_slide_text(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        text_items = self._extract_text_items(slide.shapes)
        total_text = "".join(item.get("raw_text", "") for item in text_items)

        return {
            "presentation_id": session.id,
            "slide_index": slide_index,
            "total_text_length": len(total_text),
            "item_count": len(text_items),
            "text_items": text_items,
        }

    def get_shape_details(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        shape = self._require_shape(slide, params)
        detail = self._detailed_shape_payload(shape)

        return {
            "presentation_id": session.id,
            "slide_index": slide_index,
            **detail,
        }

    # ------------------------------------------------------------------ #
    # Phase 3B — Table support
    # ------------------------------------------------------------------ #

    def add_table(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        rows = int(params["rows"])
        cols = int(params["cols"])
        ensure(rows >= 1, "validation_error", "rows must be >= 1")
        ensure(cols >= 1, "validation_error", "cols must be >= 1")

        left = to_emu(params["left"])
        top = to_emu(params["top"])
        width = to_emu(params["width"])
        height = to_emu(params["height"])

        from pptx.util import Emu as EmuVal

        table_shape = slide.shapes.add_table(rows, cols, EmuVal(left), EmuVal(top), EmuVal(width), EmuVal(height))
        table = table_shape.table

        # Optional initial data
        data = params.get("data")
        if data:
            for r_idx, row_data in enumerate(data):
                if r_idx >= rows:
                    break
                for c_idx, cell_text in enumerate(row_data):
                    if c_idx >= cols:
                        break
                    table.cell(r_idx, c_idx).text = str(cell_text)

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "shape_name": table_shape.name,
            "shape_id": int(table_shape.shape_id),
            "rows": rows,
            "cols": cols,
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def get_table(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        shape = self._find_shape(slide, params.get("shape_name") or params.get("shape_id"))
        ensure(shape.has_table, "conflict", "Shape is not a table.")

        table = shape.table
        cells = []
        for r_idx in range(len(table.rows)):
            row_cells = []
            for c_idx in range(len(table.columns)):
                cell = table.cell(r_idx, c_idx)
                cell_data = self._extract_cell_data(cell)
                row_cells.append(cell_data)
            cells.append(row_cells)

        return {
            "presentation_id": session.id,
            "slide_index": slide_index,
            "shape_name": shape.name,
            "rows": len(table.rows),
            "cols": len(table.columns),
            "cells": cells,
        }

    def set_table_cell(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        shape = self._find_shape(slide, params.get("shape_name") or params.get("shape_id"))
        ensure(shape.has_table, "conflict", "Shape is not a table.")

        table = shape.table
        row = int(params["row"])
        col = int(params["col"])
        ensure(0 <= row < len(table.rows), "validation_error", "row out of bounds")
        ensure(0 <= col < len(table.columns), "validation_error", "col out of bounds")

        cell = table.cell(row, col)
        self._write_cell(cell, params)

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "shape_name": shape.name,
            "row": row,
            "col": col,
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def set_table_data(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        shape = self._find_shape(slide, params.get("shape_name") or params.get("shape_id"))
        ensure(shape.has_table, "conflict", "Shape is not a table.")

        table = shape.table
        data = params["data"]

        for r_idx, row_data in enumerate(data):
            if r_idx >= len(table.rows):
                break
            for c_idx, cell_data in enumerate(row_data):
                if c_idx >= len(table.columns):
                    break
                cell = table.cell(r_idx, c_idx)
                if isinstance(cell_data, str):
                    cell.text = cell_data
                elif isinstance(cell_data, dict):
                    self._write_cell(cell, cell_data)

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "shape_name": shape.name,
            "rows_written": min(len(data), len(table.rows)),
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    # ------------------------------------------------------------------ #
    # Phase 3C — Shapes, notes, extras
    # ------------------------------------------------------------------ #

    _SHAPE_TYPE_MAP = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "oval": MSO_SHAPE.OVAL,
        "diamond": MSO_SHAPE.DIAMOND,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "right_arrow": MSO_SHAPE.RIGHT_ARROW,
        "left_arrow": MSO_SHAPE.LEFT_ARROW,
        "up_arrow": MSO_SHAPE.UP_ARROW,
        "down_arrow": MSO_SHAPE.DOWN_ARROW,
        "pentagon": MSO_SHAPE.REGULAR_PENTAGON,
        "hexagon": MSO_SHAPE.HEXAGON,
        "chevron": MSO_SHAPE.CHEVRON,
        "star_5_point": MSO_SHAPE.STAR_5_POINT,
        "line_inverse": MSO_SHAPE.LINE_INVERSE,
        "cross": MSO_SHAPE.CROSS,
        "frame": MSO_SHAPE.FRAME,
        "rectangular_callout": MSO_SHAPE.RECTANGULAR_CALLOUT,
        "rounded_rectangular_callout": MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
        "cloud_callout": MSO_SHAPE.CLOUD_CALLOUT,
        "cloud": MSO_SHAPE.CLOUD,
    }

    def add_shape(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        shape_type_str = str(params["shape_type"]).lower()
        mso_shape = self._SHAPE_TYPE_MAP.get(shape_type_str)
        ensure(
            mso_shape is not None,
            "validation_error",
            f"Unknown shape_type '{shape_type_str}'.",
            {"available": sorted(self._SHAPE_TYPE_MAP.keys())},
        )

        left = to_emu(params["left"])
        top = to_emu(params["top"])
        width = to_emu(params["width"])
        height = to_emu(params["height"])

        from pptx.util import Emu as EmuVal

        shape = slide.shapes.add_shape(mso_shape, EmuVal(left), EmuVal(top), EmuVal(width), EmuVal(height))

        # Optional fill color
        fill_hex = params.get("fill_hex")
        if fill_hex:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(normalize_color(str(fill_hex)))

        # Optional line color
        line_hex = params.get("line_hex")
        if line_hex:
            shape.line.color.rgb = RGBColor.from_string(normalize_color(str(line_hex)))

        # Optional text
        text = params.get("text")
        if text and hasattr(shape, "text_frame"):
            shape.text_frame.paragraphs[0].text = str(text)

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "shape_name": shape.name,
            "shape_id": int(shape.shape_id),
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def delete_shape(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        shape = self._require_shape(slide, params)
        deleted_shape = shape.name if shape.name else str(shape.shape_id)
        # Remove the shape element from the shape tree
        sp_tree = slide.shapes._spTree
        sp_tree.remove(shape.element)

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "deleted_shape": deleted_shape,
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def set_slide_notes(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        notes_text = str(params.get("notes_text", ""))
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "notes_length": len(notes_text),
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def set_shape_text(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        shape = self._require_shape(slide, params)
        ensure(
            hasattr(shape, "has_text_frame") and shape.has_text_frame,
            "conflict",
            "Shape does not have a text frame.",
        )

        text_frame = shape.text_frame
        paragraphs_data = params.get("paragraphs")

        if paragraphs_data:
            self._write_paragraphs(text_frame, paragraphs_data)
        else:
            text_content = str(params.get("text_content", ""))
            text_frame.paragraphs[0].text = text_content

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "shape_name": shape.name,
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def get_slide_xml(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        xml_bytes = etree.tostring(slide._element, pretty_print=True, encoding="unicode")

        return {
            "presentation_id": session.id,
            "slide_index": slide_index,
            "xml_content": xml_bytes,
        }

    # ------------------------------------------------------------------ #
    # Phase 4 — Flexibility tools
    # ------------------------------------------------------------------ #

    def set_shape_properties(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        shape = self._require_shape(slide, params)

        # Position / size
        if "left" in params:
            shape.left = Emu(to_emu(params["left"]))
        if "top" in params:
            shape.top = Emu(to_emu(params["top"]))
        if "width" in params:
            shape.width = Emu(to_emu(params["width"]))
        if "height" in params:
            shape.height = Emu(to_emu(params["height"]))

        # Rotation
        if "rotation" in params:
            shape.rotation = float(params["rotation"])

        # Fill
        fill_hex = params.get("fill_hex")
        if fill_hex == "none":
            shape.fill.background()
        elif fill_hex:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(normalize_color(str(fill_hex)))

        # Line / outline
        line_hex = params.get("line_hex")
        if line_hex == "none":
            shape.line.fill.background()
        elif line_hex:
            shape.line.color.rgb = RGBColor.from_string(normalize_color(str(line_hex)))

        line_width_pt = params.get("line_width_pt")
        if line_width_pt is not None:
            shape.line.width = Pt(float(line_width_pt))

        # Name
        if "name" in params:
            shape.name = str(params["name"])

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "shape_name": shape.name,
            "left_inches": emu_to_inches(shape.left) if shape.left else None,
            "top_inches": emu_to_inches(shape.top) if shape.top else None,
            "width_inches": emu_to_inches(shape.width) if shape.width else None,
            "height_inches": emu_to_inches(shape.height) if shape.height else None,
            "rotation": float(shape.rotation) if hasattr(shape, "rotation") else 0.0,
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def clone_shape(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        shape = self._require_shape(slide, params)

        # Target slide (default: same slide)
        target_slide_index = int(params.get("target_slide_index", slide_index))
        target_slide = self._get_slide_obj(session, target_slide_index)

        # Deep copy the shape XML element
        new_element = copy.deepcopy(shape.element)

        # Optional offset
        offset_left = params.get("offset_left")
        offset_top = params.get("offset_top")

        target_slide.shapes._spTree.append(new_element)

        # Find the new shape so we can apply offset
        # The cloned shape is the last one added
        new_shape = target_slide.shapes[-1] if len(target_slide.shapes) > 0 else None
        if new_shape and offset_left is not None:
            new_shape.left = Emu(int(new_shape.left or 0) + to_emu(offset_left))
        if new_shape and offset_top is not None:
            new_shape.top = Emu(int(new_shape.top or 0) + to_emu(offset_top))

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "source_slide_index": slide_index,
            "target_slide_index": target_slide_index,
            "cloned_shape_name": new_shape.name if new_shape else "unknown",
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def group_shapes(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        shape_names = params.get("shape_names", [])
        shape_ids = params.get("shape_ids", [])
        ensure(
            len(shape_names) >= 2 or len(shape_ids) >= 2,
            "validation_error",
            "At least 2 shapes required for grouping.",
        )

        identifiers = shape_names if shape_names else shape_ids
        shapes_to_group = [self._find_shape(slide, ident) for ident in identifiers]

        sp_tree = slide.shapes._spTree

        # Calculate bounding box
        min_left = min(s.left for s in shapes_to_group if s.left is not None)
        min_top = min(s.top for s in shapes_to_group if s.top is not None)
        max_right = max((s.left or 0) + (s.width or 0) for s in shapes_to_group)
        max_bottom = max((s.top or 0) + (s.height or 0) for s in shapes_to_group)

        # Create group shape element
        grp_sp = etree.SubElement(sp_tree, "{http://schemas.openxmlformats.org/presentationml/2006/main}grpSp")

        # Group shape properties
        grp_sp_pr = etree.SubElement(grp_sp, "{http://schemas.openxmlformats.org/presentationml/2006/main}grpSpPr")
        xfrm = etree.SubElement(grp_sp_pr, "{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm")
        etree.SubElement(
            xfrm,
            "{http://schemas.openxmlformats.org/drawingml/2006/main}off",
            attrib={"x": str(min_left), "y": str(min_top)},
        )
        etree.SubElement(
            xfrm,
            "{http://schemas.openxmlformats.org/drawingml/2006/main}ext",
            attrib={"cx": str(max_right - min_left), "cy": str(max_bottom - min_top)},
        )
        etree.SubElement(
            xfrm,
            "{http://schemas.openxmlformats.org/drawingml/2006/main}chOff",
            attrib={"x": str(min_left), "y": str(min_top)},
        )
        etree.SubElement(
            xfrm,
            "{http://schemas.openxmlformats.org/drawingml/2006/main}chExt",
            attrib={"cx": str(max_right - min_left), "cy": str(max_bottom - min_top)},
        )

        # Non-visual group shape properties
        next_id = max((int(s.shape_id) for s in slide.shapes), default=0) + 1
        nvGrpSpPr = etree.SubElement(grp_sp, "{http://schemas.openxmlformats.org/presentationml/2006/main}nvGrpSpPr")
        cNvPr = etree.SubElement(nvGrpSpPr, "{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr")
        cNvPr.set("id", str(next_id))
        cNvPr.set("name", str(params.get("group_name", "Group")))
        etree.SubElement(nvGrpSpPr, "{http://schemas.openxmlformats.org/presentationml/2006/main}cNvGrpSpPr")
        etree.SubElement(nvGrpSpPr, "{http://schemas.openxmlformats.org/presentationml/2006/main}nvPr")

        # Move shapes into group
        for shape in shapes_to_group:
            sp_tree.remove(shape.element)
            grp_sp.append(shape.element)

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "group_name": str(params.get("group_name", "Group")),
            "shapes_grouped": len(shapes_to_group),
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def ungroup_shapes(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        shape = self._require_shape(slide, params)

        ensure(
            shape.shape_type == MSO_SHAPE_TYPE.GROUP,
            "conflict",
            "Shape is not a group.",
        )

        sp_tree = slide.shapes._spTree
        group_element = shape.element

        # Collect children before removing group
        children = list(group_element)
        child_names: list[str] = []

        for child in children:
            # Skip non-visual properties (grpSpPr, nvGrpSpPr)
            tag = etree.QName(child).localname
            if tag in ("grpSpPr", "nvGrpSpPr"):
                continue
            group_element.remove(child)
            sp_tree.append(child)
            # Try to get name
            cNvPr = child.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr")
            if cNvPr is None:
                # Try presentation namespace
                for ns_prefix in [
                    "http://schemas.openxmlformats.org/presentationml/2006/main",
                    "http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing",
                ]:
                    cNvPr = child.find(f".//{{{ns_prefix}}}cNvPr")
                    if cNvPr is not None:
                        break
            name = cNvPr.get("name", "Unknown") if cNvPr is not None else "Unknown"
            child_names.append(name)

        # Remove the now-empty group
        sp_tree.remove(group_element)

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "ungrouped_shapes": child_names,
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def set_shape_z_order(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        shape = self._require_shape(slide, params)

        action = str(params.get("action", "front")).lower()
        ensure(
            action in ("front", "back", "forward", "backward"),
            "validation_error",
            "action must be 'front', 'back', 'forward', or 'backward'.",
        )

        sp_tree = slide.shapes._spTree
        element = shape.element

        if action == "front":
            sp_tree.remove(element)
            sp_tree.append(element)
        elif action == "back":
            sp_tree.remove(element)
            sp_tree.insert(2, element)  # index 0,1 are typically nvGrpSpPr and grpSpPr
        elif action == "forward":
            idx = list(sp_tree).index(element)
            sp_tree.remove(element)
            sp_tree.insert(min(idx + 1, len(sp_tree)), element)
        elif action == "backward":
            idx = list(sp_tree).index(element)
            sp_tree.remove(element)
            sp_tree.insert(max(idx - 1, 2), element)

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "shape_name": shape.name,
            "action": action,
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def add_image(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        image_path = str(params["image_path"])
        ensure(Path(image_path).is_file(), "not_found", f"Image file not found: {image_path}")

        left = to_emu(params.get("left", "0in"))
        top = to_emu(params.get("top", "0in"))

        # Width/height: at least one should be specified, or both
        width_param = params.get("width")
        height_param = params.get("height")

        kwargs: dict[str, Any] = {
            "image_file": image_path,
            "left": Emu(left),
            "top": Emu(top),
        }
        if width_param:
            kwargs["width"] = Emu(to_emu(width_param))
        if height_param:
            kwargs["height"] = Emu(to_emu(height_param))

        pic = slide.shapes.add_picture(**kwargs)

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "shape_name": pic.name,
            "shape_id": int(pic.shape_id),
            "left_inches": emu_to_inches(pic.left),
            "top_inches": emu_to_inches(pic.top),
            "width_inches": emu_to_inches(pic.width),
            "height_inches": emu_to_inches(pic.height),
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def add_line(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        begin_x = to_emu(params["begin_x"])
        begin_y = to_emu(params["begin_y"])
        end_x = to_emu(params["end_x"])
        end_y = to_emu(params["end_y"])

        # python-pptx add_connector doesn't exist easily, use XML approach
        from pptx.oxml.ns import qn

        sp_tree = slide.shapes._spTree

        # Calculate position and dimensions from endpoints
        left = min(begin_x, end_x)
        top = min(begin_y, end_y)
        cx = abs(end_x - begin_x)
        cy = abs(end_y - begin_y)

        # Determine if line needs flipping
        flip_h = "1" if end_x < begin_x else "0"
        flip_v = "1" if end_y < begin_y else "0"

        # Build connector shape XML
        cxn_sp = etree.SubElement(sp_tree, qn("p:cxnSp"))

        # Non-visual properties — generate a unique shape ID
        next_id = max((int(s.shape_id) for s in slide.shapes), default=0) + 1
        nv = etree.SubElement(cxn_sp, qn("p:nvCxnSpPr"))
        cNvPr = etree.SubElement(nv, qn("p:cNvPr"))
        cNvPr.set("id", str(next_id))
        cNvPr.set("name", str(params.get("line_name", "Connector")))
        etree.SubElement(nv, qn("p:cNvCxnSpPr"))
        etree.SubElement(nv, qn("p:nvPr"))

        # Shape properties
        sp_pr = etree.SubElement(cxn_sp, qn("p:spPr"))
        xfrm = etree.SubElement(sp_pr, qn("a:xfrm"))
        if flip_h == "1":
            xfrm.set("flipH", "1")
        if flip_v == "1":
            xfrm.set("flipV", "1")
        etree.SubElement(xfrm, qn("a:off"), attrib={"x": str(left), "y": str(top)})
        etree.SubElement(xfrm, qn("a:ext"), attrib={"cx": str(max(cx, 1)), "cy": str(max(cy, 1))})

        # Preset geometry (straight line)
        prst_geom = etree.SubElement(sp_pr, qn("a:prstGeom"), attrib={"prst": "line"})
        etree.SubElement(prst_geom, qn("a:avLst"))

        # Line properties
        line_hex = params.get("color_hex")
        line_width_pt = params.get("width_pt", 1.0)

        ln = etree.SubElement(sp_pr, qn("a:ln"), attrib={"w": str(int(Pt(float(line_width_pt))))})
        if line_hex:
            solid = etree.SubElement(ln, qn("a:solidFill"))
            etree.SubElement(solid, qn("a:srgbClr"), attrib={"val": normalize_color(str(line_hex))})
        else:
            solid = etree.SubElement(ln, qn("a:solidFill"))
            etree.SubElement(solid, qn("a:srgbClr"), attrib={"val": "000000"})

        # Dash style
        dash_style = params.get("dash_style")
        if dash_style:
            etree.SubElement(ln, qn("a:prstDash"), attrib={"val": str(dash_style)})

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "line_name": str(params.get("line_name", "Connector")),
            "begin": {"x_inches": emu_to_inches(begin_x), "y_inches": emu_to_inches(begin_y)},
            "end": {"x_inches": emu_to_inches(end_x), "y_inches": emu_to_inches(end_y)},
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def find_replace_text(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        prs = self._prs(session)

        find_text = str(params["find_text"])
        replace_text = str(params["replace_text"])
        case_sensitive = bool(params.get("case_sensitive", True))

        # Optional: limit to specific slides
        slide_indices = params.get("slide_indices")
        if slide_indices:
            slides = [prs.slides[i - 1] for i in slide_indices if 1 <= i <= len(prs.slides)]
        else:
            slides = list(prs.slides)

        total_replacements = 0

        for slide in slides:
            total_replacements += self._replace_in_shapes(slide.shapes, find_text, replace_text, case_sensitive)

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "find_text": find_text,
            "replace_text": replace_text,
            "total_replacements": total_replacements,
            "slides_searched": len(slides),
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    # ------------------------------------------------------------------ #
    # Phase 5 — Chart tools
    # ------------------------------------------------------------------ #

    _CHART_TYPE_MAP = {
        "column_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
        "column_stacked_100": XL_CHART_TYPE.COLUMN_STACKED_100,
        "bar_clustered": XL_CHART_TYPE.BAR_CLUSTERED,
        "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
        "bar_stacked_100": XL_CHART_TYPE.BAR_STACKED_100,
        "line": XL_CHART_TYPE.LINE,
        "line_markers": XL_CHART_TYPE.LINE_MARKERS,
        "line_stacked": XL_CHART_TYPE.LINE_STACKED,
        "pie": XL_CHART_TYPE.PIE,
        "pie_exploded": XL_CHART_TYPE.PIE_EXPLODED,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
        "area": XL_CHART_TYPE.AREA,
        "area_stacked": XL_CHART_TYPE.AREA_STACKED,
        "area_stacked_100": XL_CHART_TYPE.AREA_STACKED_100,
        "xy_scatter": XL_CHART_TYPE.XY_SCATTER,
        "xy_scatter_lines": XL_CHART_TYPE.XY_SCATTER_LINES,
        "xy_scatter_smooth": XL_CHART_TYPE.XY_SCATTER_SMOOTH,
        "bubble": XL_CHART_TYPE.BUBBLE,
        "radar": XL_CHART_TYPE.RADAR,
        "stock_hlc": XL_CHART_TYPE.STOCK_HLC,
        "stock_ohlc": XL_CHART_TYPE.STOCK_OHLC,
        "three_d_column": XL_CHART_TYPE.THREE_D_COLUMN,
        "three_d_bar_clustered": XL_CHART_TYPE.THREE_D_BAR_CLUSTERED,
        "three_d_pie": XL_CHART_TYPE.THREE_D_PIE,
        "three_d_line": XL_CHART_TYPE.THREE_D_LINE,
    }

    _LEGEND_POSITION_MAP = {
        "bottom": XL_LEGEND_POSITION.BOTTOM,
        "corner": XL_LEGEND_POSITION.CORNER,
        "left": XL_LEGEND_POSITION.LEFT,
        "right": XL_LEGEND_POSITION.RIGHT,
        "top": XL_LEGEND_POSITION.TOP,
    }

    def add_chart(self, params: dict[str, Any]) -> dict[str, Any]:
        from pptx.chart.data import BubbleChartData, CategoryChartData, XyChartData

        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        chart_type_str = str(params["chart_type"]).lower()
        xl_type = self._CHART_TYPE_MAP.get(chart_type_str)
        ensure(
            xl_type is not None,
            "validation_error",
            f"Unknown chart_type '{chart_type_str}'.",
            {"available": sorted(self._CHART_TYPE_MAP.keys())},
        )

        left = to_emu(params.get("left", "1in"))
        top = to_emu(params.get("top", "1in"))
        width = to_emu(params.get("width", "8in"))
        height = to_emu(params.get("height", "5in"))

        series_data = params.get("series", [])
        ensure(len(series_data) > 0, "validation_error", "At least one series is required.")

        # Determine data class based on chart type
        is_xy = chart_type_str.startswith("xy_scatter")
        is_bubble = chart_type_str.startswith("bubble")

        if is_xy:
            chart_data = XyChartData()
            for s in series_data:
                series_obj = chart_data.add_series(str(s.get("name", "Series")))
                data_points = s.get("data_points", [])
                for dp in data_points:
                    series_obj.add_data_point(float(dp["x"]), float(dp["y"]))
        elif is_bubble:
            chart_data = BubbleChartData()
            for s in series_data:
                series_obj = chart_data.add_series(str(s.get("name", "Series")))
                data_points = s.get("data_points", [])
                for dp in data_points:
                    series_obj.add_data_point(float(dp["x"]), float(dp["y"]), float(dp.get("size", 10)))
        else:
            chart_data = CategoryChartData()
            categories = params.get("categories", [])
            if categories:
                chart_data.categories = categories
            for s in series_data:
                chart_data.add_series(str(s.get("name", "Series")), s.get("values", []))

        graphic_frame = slide.shapes.add_chart(xl_type, Emu(left), Emu(top), Emu(width), Emu(height), chart_data)
        chart = graphic_frame.chart

        # Optional styling
        if params.get("has_legend") is not None:
            chart.has_legend = bool(params["has_legend"])

        if chart.has_legend and params.get("legend_position"):
            pos = self._LEGEND_POSITION_MAP.get(str(params["legend_position"]).lower())
            if pos is not None:
                chart.legend.position = pos
                chart.legend.include_in_layout = False

        if params.get("has_data_labels"):
            plot = chart.plots[0]
            plot.has_data_labels = True
            fmt = params.get("data_label_number_format")
            if fmt:
                plot.data_labels.number_format = str(fmt)

        if params.get("chart_style") is not None:
            chart.style = int(params["chart_style"])

        if params.get("title"):
            chart.has_title = True
            chart.chart_title.text_frame.text = str(params["title"])

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "shape_name": graphic_frame.name,
            "shape_id": int(graphic_frame.shape_id),
            "chart_type": chart_type_str,
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def get_chart_data(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        shape = self._require_shape(slide, params)

        ensure(shape.has_chart, "conflict", "Shape does not contain a chart.")
        chart = shape.chart

        # Determine chart type name
        chart_type_name = "unknown"
        for name, xl_type in self._CHART_TYPE_MAP.items():
            if chart.chart_type == xl_type:
                chart_type_name = name
                break

        # Extract categories
        categories: list[str] = []
        try:
            plot = chart.plots[0]
            for cat in plot.categories:
                categories.append(str(cat) if cat is not None else "")
        except Exception as exc:
            print(f"[pptx] Warning: could not extract chart categories: {exc}", file=sys.stderr)

        # Extract series
        series_list: list[dict[str, Any]] = []
        try:
            for series in chart.series:
                s_data: dict[str, Any] = {"name": str(series.name) if hasattr(series, "name") else ""}
                # Try to get values
                try:
                    s_data["values"] = [float(v) if v is not None else None for v in series.values]
                except Exception as exc:
                    s_data["values"] = []
                    print(f"[pptx] Warning: could not extract series values: {exc}", file=sys.stderr)
                series_list.append(s_data)
        except Exception as exc:
            print(f"[pptx] Warning: could not extract chart series: {exc}", file=sys.stderr)

        # Chart properties
        result: dict[str, Any] = {
            "presentation_id": session.id,
            "slide_index": slide_index,
            "shape_name": shape.name,
            "chart_type": chart_type_name,
            "categories": categories,
            "series": series_list,
            "has_legend": bool(chart.has_legend),
            "has_title": bool(chart.has_title),
        }

        if chart.has_title:
            try:
                result["title"] = chart.chart_title.text_frame.text
            except Exception:
                result["title"] = ""

        # Data labels
        try:
            plot = chart.plots[0]
            result["has_data_labels"] = bool(plot.has_data_labels)
            if plot.has_data_labels:
                result["data_label_number_format"] = str(plot.data_labels.number_format)
        except Exception:
            result["has_data_labels"] = False

        return result

    def update_chart_data(self, params: dict[str, Any]) -> dict[str, Any]:
        from pptx.chart.data import CategoryChartData, XyChartData

        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        shape = self._require_shape(slide, params)

        ensure(shape.has_chart, "conflict", "Shape does not contain a chart.")
        chart = shape.chart

        # Determine if XY/scatter type
        chart_type_name = "unknown"
        for name, xl_type in self._CHART_TYPE_MAP.items():
            if chart.chart_type == xl_type:
                chart_type_name = name
                break

        is_xy = chart_type_name.startswith("xy_scatter")
        series_data = params.get("series", [])

        if is_xy:
            new_data = XyChartData()
            for s in series_data:
                series_obj = new_data.add_series(str(s.get("name", "Series")))
                for dp in s.get("data_points", []):
                    series_obj.add_data_point(float(dp["x"]), float(dp["y"]))
        else:
            new_data = CategoryChartData()
            categories = params.get("categories")
            if categories:
                new_data.categories = categories
            else:
                # Keep existing categories
                try:
                    existing_cats = [str(c) for c in chart.plots[0].categories]
                    new_data.categories = existing_cats
                except Exception:
                    pass

            for s in series_data:
                new_data.add_series(str(s.get("name", "Series")), s.get("values", []))

        chart.replace_data(new_data)
        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "shape_name": shape.name,
            "chart_type": chart_type_name,
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def set_chart_style(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        shape = self._require_shape(slide, params)

        ensure(shape.has_chart, "conflict", "Shape does not contain a chart.")
        chart = shape.chart

        # Legend
        if "has_legend" in params:
            chart.has_legend = bool(params["has_legend"])

        if chart.has_legend and params.get("legend_position"):
            pos = self._LEGEND_POSITION_MAP.get(str(params["legend_position"]).lower())
            if pos is not None:
                chart.legend.position = pos
                chart.legend.include_in_layout = bool(params.get("legend_in_layout", False))

        # Data labels
        if "has_data_labels" in params:
            plot = chart.plots[0]
            plot.has_data_labels = bool(params["has_data_labels"])
            if plot.has_data_labels:
                fmt = params.get("data_label_number_format")
                if fmt:
                    plot.data_labels.number_format = str(fmt)

                dl_position = params.get("data_label_position")
                if dl_position:
                    from pptx.enum.chart import XL_LABEL_POSITION

                    dl_pos_map = {
                        "center": XL_LABEL_POSITION.CENTER,
                        "inside_end": XL_LABEL_POSITION.INSIDE_END,
                        "outside_end": XL_LABEL_POSITION.OUTSIDE_END,
                        "inside_base": XL_LABEL_POSITION.INSIDE_BASE,
                        "above": XL_LABEL_POSITION.ABOVE,
                        "below": XL_LABEL_POSITION.BELOW,
                        "left": XL_LABEL_POSITION.LEFT,
                        "right": XL_LABEL_POSITION.RIGHT,
                        "best_fit": XL_LABEL_POSITION.BEST_FIT,
                    }
                    pos = dl_pos_map.get(str(dl_position).lower())
                    if pos is not None:
                        plot.data_labels.position = pos

        # Chart style
        if "chart_style" in params:
            chart.style = int(params["chart_style"])

        # Title
        if "title" in params:
            title_text = params["title"]
            if title_text is None or title_text == "":
                chart.has_title = False
            else:
                chart.has_title = True
                chart.chart_title.text_frame.text = str(title_text)

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "shape_name": shape.name,
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    # ------------------------------------------------------------------ #
    # Phase 6 — Agent workflow tools
    # ------------------------------------------------------------------ #

    def copy_shape_between_decks(self, params: dict[str, Any]) -> dict[str, Any]:
        source_session = self._get_session(str(params["source_presentation_id"]))
        target_session = self._get_session(str(params["target_presentation_id"]))

        source_slide_index = int(params["source_slide_index"])
        target_slide_index = int(params["target_slide_index"])

        source_slide = self._get_slide_obj(source_session, source_slide_index)
        target_slide = self._get_slide_obj(target_session, target_slide_index)

        source_shape = self._require_shape(source_slide, params)

        # Deep copy the shape XML element
        new_element = copy.deepcopy(source_shape.element)

        # Handle image relationships: copy image parts from source to target
        source_slide_part = source_slide.part
        target_slide_part = target_slide.part

        # Find all relationship references (r:blipFill, r:embed, etc.)
        nsmap = {
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        }
        # Transfer blip (image) relationships
        for blip in new_element.findall(".//a:blip", nsmap):
            r_embed = blip.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
            if r_embed:
                try:
                    source_rel = source_slide_part.rels[r_embed]
                    image_part = source_rel.target_part
                    # Add the image part to the target slide and get new rId
                    new_rel = target_slide_part.relate_to(image_part, source_rel.reltype)
                    blip.set("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed", new_rel)
                except Exception as exc:
                    print(f"[pptx] Warning: could not transfer image relationship {r_embed}: {exc}", file=sys.stderr)

        # Apply offset if specified
        offset_left = params.get("offset_left")
        offset_top = params.get("offset_top")
        if offset_left is not None:
            try:
                off_elem = new_element.find(".//{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetml}off")
                if off_elem is None:
                    # Try presentation ML namespace
                    for off in new_element.iter():
                        if off.tag.endswith("}off") or off.tag == "off":
                            off_elem = off
                            break
                if off_elem is not None and off_elem.get("x") is not None:
                    current_x = int(off_elem.get("x", "0"))
                    off_elem.set("x", str(current_x + to_emu(offset_left)))
            except Exception as exc:
                print(f"[pptx] Warning: could not apply offset_left: {exc}", file=sys.stderr)

        if offset_top is not None:
            try:
                for off in new_element.iter():
                    if (off.tag.endswith("}off") or off.tag == "off") and off.get("y") is not None:
                        current_y = int(off.get("y", "0"))
                        off.set("y", str(current_y + to_emu(offset_top)))
                        break
            except Exception as exc:
                print(f"[pptx] Warning: could not apply offset_top: {exc}", file=sys.stderr)

        # Insert into target slide's shape tree
        target_sp_tree = target_slide.shapes._spTree
        target_sp_tree.append(new_element)

        self._persist(target_session)

        # Get the new shape's name from the inserted element
        new_shape_name = ""
        try:
            nv_elem = new_element.find(".//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr")
            if nv_elem is None:
                nv_elem = new_element.find(".//{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetml}cNvPr")
            if nv_elem is None:
                # Generic search for cNvPr
                for elem in new_element.iter():
                    if elem.tag.endswith("}cNvPr") or elem.tag == "cNvPr":
                        nv_elem = elem
                        break
            if nv_elem is not None:
                new_shape_name = nv_elem.get("name", "")
        except Exception:
            pass

        return {
            "success": True,
            "source_presentation_id": source_session.id,
            "target_presentation_id": target_session.id,
            "source_slide_index": source_slide_index,
            "target_slide_index": target_slide_index,
            "source_shape_name": source_shape.name,
            "new_shape_name": new_shape_name,
            "presentation_state": self.get_presentation_state({"presentation_id": target_session.id}),
        }

    def get_slide_shapes(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        shapes_list: list[dict[str, Any]] = []
        for shape in slide.shapes:
            shape_info: dict[str, Any] = {
                "shape_id": int(shape.shape_id),
                "name": shape.name,
                "shape_type": str(shape.shape_type).split("(")[0].strip() if shape.shape_type else "unknown",
                "left": emu_to_inches(shape.left) if shape.left is not None else None,
                "top": emu_to_inches(shape.top) if shape.top is not None else None,
                "width": emu_to_inches(shape.width) if shape.width is not None else None,
                "height": emu_to_inches(shape.height) if shape.height is not None else None,
            }

            # Add key flags
            shape_info["is_placeholder"] = shape.is_placeholder
            shape_info["has_text_frame"] = shape.has_text_frame
            shape_info["has_table"] = shape.has_table
            shape_info["has_chart"] = shape.has_chart

            if shape.is_placeholder:
                try:
                    shape_info["placeholder_name"] = shape.placeholder_format.idx
                except Exception:
                    pass

            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                shape_info["child_count"] = len(shape.shapes)

            shapes_list.append(shape_info)

        return {
            "presentation_id": session.id,
            "slide_index": slide_index,
            "shape_count": len(shapes_list),
            "shapes": shapes_list,
        }

    def set_table_cell_merge(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        shape = self._require_shape(slide, params)

        ensure(shape.has_table, "conflict", "Shape does not contain a table.")
        table = shape.table

        start_row = int(params["start_row"])
        start_col = int(params["start_col"])
        end_row = int(params["end_row"])
        end_col = int(params["end_col"])

        ensure(
            0 <= start_row <= end_row < len(table.rows),
            "validation_error",
            f"Row indices must satisfy 0 <= start_row <= end_row < {len(table.rows)}.",
        )
        ensure(
            0 <= start_col <= end_col < len(table.columns),
            "validation_error",
            f"Col indices must satisfy 0 <= start_col <= end_col < {len(table.columns)}.",
        )

        # Use python-pptx's native merge
        start_cell = table.cell(start_row, start_col)
        end_cell = table.cell(end_row, end_col)
        start_cell.merge(end_cell)

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "shape_name": shape.name,
            "merged_range": f"({start_row},{start_col})->({end_row},{end_col})",
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    # ------------------------------------------------------------------ #
    # Phase 7 — Formatting & Fidelity tools
    # ------------------------------------------------------------------ #

    def set_paragraph_spacing(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        shape = self._require_shape(slide, params)
        ensure(shape.has_text_frame, "conflict", "Shape does not contain a text frame.")

        paragraph_index = int(params["paragraph_index"])
        ensure(
            0 <= paragraph_index < len(shape.text_frame.paragraphs),
            "validation_error",
            f"Paragraph index out of bounds. Max: {len(shape.text_frame.paragraphs) - 1}",
        )
        paragraph = shape.text_frame.paragraphs[paragraph_index]

        if "line_spacing" in params:
            val = params["line_spacing"]
            paragraph.line_spacing = Pt(float(val)) if val is not None else None

        if "space_before" in params:
            val = params["space_before"]
            paragraph.space_before = Pt(float(val)) if val is not None else None

        if "space_after" in params:
            val = params["space_after"]
            paragraph.space_after = Pt(float(val)) if val is not None else None

        self._persist(session)

        return {
            "success": True,
            "presentation_id": session.id,
            "slide_index": slide_index,
            "shape_name": shape.name,
            "paragraph_index": paragraph_index,
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def set_text_box_properties(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        shape = self._require_shape(slide, params)
        ensure(shape.has_text_frame, "conflict", "Shape does not contain a text frame.")

        tf = shape.text_frame

        if "margin_left" in params:
            tf.margin_left = to_emu(params["margin_left"]) if params["margin_left"] else 0
        if "margin_right" in params:
            tf.margin_right = to_emu(params["margin_right"]) if params["margin_right"] else 0
        if "margin_top" in params:
            tf.margin_top = to_emu(params["margin_top"]) if params["margin_top"] else 0
        if "margin_bottom" in params:
            tf.margin_bottom = to_emu(params["margin_bottom"]) if params["margin_bottom"] else 0

        if "word_wrap" in params:
            tf.word_wrap = bool(params["word_wrap"]) if params["word_wrap"] is not None else None

        if "auto_size" in params:
            val = params["auto_size"]
            if val is None:
                tf.auto_size = None
            else:
                from pptx.enum.text import MSO_AUTO_SIZE

                auto_map = {
                    "none": MSO_AUTO_SIZE.NONE,
                    "shape_to_fit_text": MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
                    "text_to_fit_shape": MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
                }
                tf.auto_size = auto_map.get(str(val).lower())

        if "vertical_alignment" in params:
            val = params["vertical_alignment"]
            if val is None:
                tf.vertical_anchor = None
            else:
                from pptx.enum.text import MSO_VERTICAL_ANCHOR

                val_map = {
                    "top": MSO_VERTICAL_ANCHOR.TOP,
                    "middle": MSO_VERTICAL_ANCHOR.MIDDLE,
                    "bottom": MSO_VERTICAL_ANCHOR.BOTTOM,
                }
                tf.vertical_anchor = val_map.get(str(val).lower())

        self._persist(session)
        return {
            "success": True,
            "shape_name": shape.name,
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def set_table_style(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        shape = self._require_shape(slide, params)
        ensure(shape.has_table, "conflict", "Shape does not contain a table.")

        table = shape.table

        if "first_row" in params:
            table.first_row = bool(params["first_row"])
        if "last_row" in params:
            table.last_row = bool(params["last_row"])
        if "first_col" in params:
            table.first_col = bool(params["first_col"])
        if "last_col" in params:
            table.last_col = bool(params["last_col"])
        if "banded_rows" in params:
            table.horz_banding = bool(params["banded_rows"])
        if "banded_cols" in params:
            table.vert_banding = bool(params["banded_cols"])

        # Hack to set style ID if provided (python-pptx doesn't formally expose style_id assignment easily)
        if "style_id" in params:
            style_id = str(params["style_id"]).strip()
            # The tblPr element holds the style_id reference
            tbl_pr = shape.element.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}tblPr")
            if tbl_pr is not None:
                tbl_style_id = tbl_pr.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}tableStyleId")
                if tbl_style_id is None:
                    # Create it if it doesn't exist
                    tbl_style_id = etree.SubElement(
                        tbl_pr, "{http://schemas.openxmlformats.org/drawingml/2006/main}tableStyleId"
                    )
                tbl_style_id.text = style_id

        self._persist(session)

        return {
            "success": True,
            "shape_name": shape.name,
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def set_shape_fill_gradient(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        shape = self._require_shape(slide, params)

        fill = shape.fill
        fill.gradient()

        if "angle" in params:
            fill.gradient_angle = float(params["angle"])

        stops = params.get("stops", [])
        if stops:
            # Clear existing stops (python-pptx usually creates 2 default stops)
            while len(fill.gradient_stops) > 0:
                # Remove from underlying XML to clear it
                gs_lst = shape.element.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}gsLst")
                if gs_lst is not None:
                    for gs in gs_lst:
                        gs_lst.remove(gs)
                # python-pptx doesn't have a direct clear() method for gradient_stops
                break

            # Adding new stops requires XML manipulation because python-pptx has limited write support for stops
            gs_lst = shape.element.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}gsLst")
            if gs_lst is None:
                grad_fill = shape.element.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}gradFill")
                if grad_fill is not None:
                    gs_lst = etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}gsLst")
                    grad_fill.insert(0, gs_lst)

            if gs_lst is not None:
                for stop in stops:
                    pos = int(float(stop.get("position", 0.0)) * 100000)  # 0 to 100000 (0% to 100%)
                    hex_color = str(stop.get("color_hex", "FFFFFF")).strip("#")

                    gs_elem = etree.SubElement(gs_lst, "{http://schemas.openxmlformats.org/drawingml/2006/main}gs")
                    gs_elem.set("pos", str(pos))

                    srgb_clr = etree.SubElement(
                        gs_elem, "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr"
                    )
                    srgb_clr.set("val", hex_color)

        self._persist(session)
        return {
            "success": True,
            "shape_name": shape.name,
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def add_connector(self, params: dict[str, Any]) -> dict[str, Any]:
        session = self._get_session(str(params["presentation_id"]))
        slide_index = int(params["slide_index"])
        slide = self._get_slide_obj(session, slide_index)

        from pptx.enum.shapes import MSO_CONNECTOR

        type_map = {
            "straight": MSO_CONNECTOR.STRAIGHT,
            "elbow": MSO_CONNECTOR.ELBOW,
            "curve": MSO_CONNECTOR.CURVE,
        }
        ctype_str = str(params.get("connector_type", "straight")).lower()
        connector_type = type_map.get(ctype_str, MSO_CONNECTOR.STRAIGHT)

        # Temporary coordinates for creation
        connector = slide.shapes.add_connector(connector_type, Emu(0), Emu(0), Emu(100), Emu(100))

        if "begin_shape_name" in params or "begin_shape_id" in params:
            b_shape = self._find_shape(slide, params.get("begin_shape_name") or params.get("begin_shape_id"))
            b_idx = int(params.get("begin_connection_site", 0))
            connector.begin_connect(b_shape, b_idx)

        if "end_shape_name" in params or "end_shape_id" in params:
            e_shape = self._find_shape(slide, params.get("end_shape_name") or params.get("end_shape_id"))
            e_idx = int(params.get("end_connection_site", 0))
            connector.end_connect(e_shape, e_idx)

        # Apply formatting
        if "color_hex" in params:
            connector.line.color.rgb = RGBColor.from_string(str(params["color_hex"]).strip("#"))

        if "width_pt" in params:
            connector.line.width = Pt(float(params["width_pt"]))

        self._persist(session)
        return {
            "success": True,
            "shape_name": connector.name,
            "shape_id": int(connector.shape_id),
            "presentation_state": self.get_presentation_state({"presentation_id": session.id}),
        }

    def shutdown(self) -> None:
        for session in list(self.sessions.values()):
            try:
                Path(session.working_path).unlink(missing_ok=True)
            except Exception:
                pass
        self.sessions.clear()

    def _get_session(self, presentation_id: str) -> EngineSession:
        session = self.sessions.get(presentation_id)
        if not session:
            raise BridgeError(
                code="not_found",
                message="Presentation session not found.",
                details={"presentation_id": presentation_id},
            )
        return session

    def _prs(self, session: EngineSession) -> Presentation:
        prs = session.extra.get("prs")
        if prs is None:
            raise BridgeError(code="engine_error", message="Session presentation object is missing.")
        return prs

    def _persist(self, session: EngineSession) -> None:
        prs = self._prs(session)
        prs.save(session.working_path)
        session.dirty = True

    def _assert_slide_index(self, prs: Presentation, slide_index: int) -> None:
        if slide_index < 1 or slide_index > len(prs.slides):
            raise BridgeError(
                code="validation_error",
                message="slide_index is out of bounds.",
                details={"slide_index": slide_index, "slide_count": len(prs.slides)},
            )

    def _get_slide_obj(self, session: EngineSession, slide_index: int):
        prs = self._prs(session)
        self._assert_slide_index(prs, slide_index)
        return prs.slides[slide_index - 1]

    def _slide_title(self, slide) -> str:
        title_shape = slide.shapes.title
        if title_shape is not None and hasattr(title_shape, "text") and title_shape.text:
            return title_shape.text.strip()

        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame and shape.text:
                return shape.text.strip().splitlines()[0]

        return ""

    def _find_layout(self, prs: Presentation, layout_name: str):
        for layout in prs.slide_layouts:
            if layout.name == layout_name:
                return layout

        available = [layout.name for layout in prs.slide_layouts]
        raise BridgeError(
            code="validation_error",
            message=f"Layout '{layout_name}' not found.",
            details={"layout_name": layout_name, "available_layouts": available},
        )

    def _placeholder_payload(self, placeholder) -> dict[str, Any]:
        return {
            "name": placeholder.name,
            "idx": int(placeholder.placeholder_format.idx),
            "type": str(placeholder.placeholder_format.type),
            "left_inches": emu_to_inches(placeholder.left),
            "top_inches": emu_to_inches(placeholder.top),
            "width_inches": emu_to_inches(placeholder.width),
            "height_inches": emu_to_inches(placeholder.height),
        }

    def _shape_payload(self, shape) -> dict[str, Any]:
        payload: dict[str, Any] = {
            "shape_id": int(shape.shape_id),
            "name": shape.name,
            "type": str(shape.shape_type),
            "left_inches": emu_to_inches(shape.left),
            "top_inches": emu_to_inches(shape.top),
            "width_inches": emu_to_inches(shape.width),
            "height_inches": emu_to_inches(shape.height),
            "is_placeholder": bool(getattr(shape, "is_placeholder", False)),
        }

        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            payload["text"] = shape.text

        return payload

    def _find_placeholder(self, slide, placeholder_name: str):
        for placeholder in slide.placeholders:
            if placeholder.name == placeholder_name:
                return placeholder

        available = [placeholder.name for placeholder in slide.placeholders]
        raise BridgeError(
            code="not_found",
            message=f"Placeholder '{placeholder_name}' not found.",
            details={"placeholder_name": placeholder_name, "available_placeholders": available},
        )

    def _validate_order(self, prs: Presentation, new_order: list[int]) -> None:
        count = len(prs.slides)
        if len(new_order) != count:
            raise BridgeError(
                code="validation_error",
                message="new_order length must equal slide_count.",
                details={"new_order_length": len(new_order), "slide_count": count},
            )

        expected = set(range(1, count + 1))
        received = set(new_order)
        if expected != received:
            raise BridgeError(
                code="validation_error",
                message="new_order must contain each current slide index exactly once.",
                details={"expected": sorted(expected), "received": new_order},
            )

    def _reorder_slide_ids(self, prs: Presentation, new_order: list[int]) -> None:
        sld_id_lst = prs.slides._sldIdLst
        slide_ids = list(sld_id_lst)
        reordered = [slide_ids[index - 1] for index in new_order]

        for node in list(sld_id_lst):
            sld_id_lst.remove(node)
        for node in reordered:
            sld_id_lst.append(node)

    def _send_shape_to_back(self, slide, shape) -> None:
        sp_tree = slide.shapes._spTree
        element = shape.element
        sp_tree.remove(element)
        sp_tree.insert(2, element)

    def _extract_theme(self, path: Path) -> dict[str, Any]:
        if not path.exists():
            return {"colors": {}, "fonts": {"major": "", "minor": ""}}

        colors: dict[str, str] = {}
        fonts = {"major": "", "minor": ""}

        with zipfile.ZipFile(path, "r") as archive:
            try:
                xml_data = archive.read("ppt/theme/theme1.xml")
            except KeyError:
                return {"colors": colors, "fonts": fonts}

        root = etree.fromstring(xml_data)

        clr_scheme = root.find(".//a:clrScheme", namespaces=_COLOR_NAMESPACES)
        if clr_scheme is not None:
            for child in list(clr_scheme):
                name = etree.QName(child).localname
                srgb = child.find(".//a:srgbClr", namespaces=_COLOR_NAMESPACES)
                sys_clr = child.find(".//a:sysClr", namespaces=_COLOR_NAMESPACES)
                if srgb is not None:
                    colors[name] = srgb.attrib.get("val", "")
                elif sys_clr is not None:
                    colors[name] = sys_clr.attrib.get("lastClr", "")

        major_font = root.find(".//a:fontScheme/a:majorFont/a:latin", namespaces=_COLOR_NAMESPACES)
        minor_font = root.find(".//a:fontScheme/a:minorFont/a:latin", namespaces=_COLOR_NAMESPACES)
        if major_font is not None:
            fonts["major"] = major_font.attrib.get("typeface", "")
        if minor_font is not None:
            fonts["minor"] = minor_font.attrib.get("typeface", "")

        return {"colors": colors, "fonts": fonts}

    # ------------------------------------------------------------------ #
    # Phase 3 helpers
    # ------------------------------------------------------------------ #

    def _apply_run_formatting(self, run, data: dict[str, Any]) -> None:
        """Apply font formatting from a dict of optional formatting keys."""
        if data.get("font_name"):
            run.font.name = str(data["font_name"])
        if data.get("font_size_pt"):
            run.font.size = Pt(float(data["font_size_pt"]))
        if data.get("bold") is not None:
            run.font.bold = bool(data["bold"])
        if data.get("italic") is not None:
            run.font.italic = bool(data["italic"])
        if data.get("underline") is not None:
            run.font.underline = bool(data["underline"])
        if data.get("color_hex"):
            run.font.color.rgb = RGBColor.from_string(normalize_color(str(data["color_hex"])))

    def _write_paragraphs(self, text_frame, paragraphs_data: list[dict[str, Any]]) -> None:
        """Write rich paragraph data into a text frame (clears existing content)."""
        text_frame.clear()
        for p_idx, p_data in enumerate(paragraphs_data):
            paragraph = text_frame.paragraphs[0] if p_idx == 0 else text_frame.add_paragraph()

            alignment = p_data.get("alignment")
            if alignment and alignment in self._ALIGN_MAP:
                paragraph.alignment = self._ALIGN_MAP[alignment]

            level = p_data.get("level")
            if level is not None:
                paragraph.level = int(level)

            runs_data = p_data.get("runs", [])
            if not runs_data:
                run = paragraph.add_run()
                run.text = str(p_data.get("text", ""))
            else:
                for r_data in runs_data:
                    run = paragraph.add_run()
                    run.text = str(r_data.get("text", ""))
                    self._apply_run_formatting(run, r_data)

    def _require_shape(self, slide, params: dict[str, Any]):
        """Resolve a shape from params by shape_name or shape_id, raising on missing."""
        identifier = params.get("shape_name") or params.get("shape_id")
        ensure(identifier is not None, "validation_error", "shape_name or shape_id is required.")
        return self._find_shape(slide, identifier)

    def _extract_text_items(self, shapes, depth: int = 0) -> list[dict[str, Any]]:
        """Recursively extract text from all shapes including tables, groups, etc."""
        items: list[dict[str, Any]] = []

        for shape in shapes:
            shape_info: dict[str, Any] = {
                "shape_name": shape.name,
                "shape_id": int(shape.shape_id),
                "shape_type": str(shape.shape_type),
                "is_placeholder": bool(getattr(shape, "is_placeholder", False)),
                "left_inches": emu_to_inches(shape.left) if shape.left is not None else None,
                "top_inches": emu_to_inches(shape.top) if shape.top is not None else None,
                "width_inches": emu_to_inches(shape.width) if shape.width is not None else None,
                "height_inches": emu_to_inches(shape.height) if shape.height is not None else None,
            }

            # --- Table shapes ---
            if shape.has_table:
                table = shape.table
                table_text_parts: list[str] = []
                table_cells: list[dict[str, Any]] = []
                for r_idx in range(len(table.rows)):
                    for c_idx in range(len(table.columns)):
                        cell = table.cell(r_idx, c_idx)
                        cell_text = cell.text or ""
                        table_text_parts.append(cell_text)
                        paragraphs = self._extract_paragraphs(cell.text_frame) if cell.text_frame else []
                        table_cells.append(
                            {
                                "row": r_idx,
                                "col": c_idx,
                                "text": cell_text,
                                "paragraphs": paragraphs,
                            }
                        )

                shape_info["content_type"] = "table"
                shape_info["rows"] = len(table.rows)
                shape_info["cols"] = len(table.columns)
                shape_info["table_cells"] = table_cells
                shape_info["raw_text"] = "\n".join(table_text_parts)
                shape_info["paragraphs"] = []
                items.append(shape_info)

            # --- Group shapes —- recurse into children ---
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                shape_info["content_type"] = "group"
                shape_info["raw_text"] = ""
                shape_info["paragraphs"] = []
                items.append(shape_info)
                # Recursively extract text from group children
                if hasattr(shape, "shapes"):
                    child_items = self._extract_text_items(shape.shapes, depth + 1)
                    for child in child_items:
                        child["parent_group"] = shape.name
                    items.extend(child_items)

            # --- Shapes with text frames (text boxes, autoshapes, placeholders) ---
            elif hasattr(shape, "has_text_frame") and shape.has_text_frame:
                paragraphs = self._extract_paragraphs(shape.text_frame)
                raw_text = shape.text or ""
                shape_info["content_type"] = "text"
                shape_info["raw_text"] = raw_text
                shape_info["paragraphs"] = paragraphs
                items.append(shape_info)

            # --- Picture shapes ---
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_info["content_type"] = "picture"
                shape_info["raw_text"] = ""
                shape_info["paragraphs"] = []
                items.append(shape_info)

            # --- Other shapes (connectors, charts, etc) ---
            else:
                shape_info["content_type"] = "other"
                shape_info["raw_text"] = ""
                shape_info["paragraphs"] = []
                items.append(shape_info)

        return items

    def _extract_paragraphs(self, text_frame) -> list[dict[str, Any]]:
        """Extract paragraph + run structure from a text frame."""
        paragraphs: list[dict[str, Any]] = []
        for paragraph in text_frame.paragraphs:
            runs: list[dict[str, Any]] = []
            for run in paragraph.runs:
                color_hex = None
                try:
                    if run.font.color and run.font.color.rgb:
                        color_hex = str(run.font.color.rgb)
                except Exception:
                    pass

                runs.append(
                    {
                        "text": run.text,
                        "bold": bool(run.font.bold) if run.font.bold is not None else None,
                        "italic": bool(run.font.italic) if run.font.italic is not None else None,
                        "underline": bool(run.font.underline) if run.font.underline is not None else None,
                        "font_name": run.font.name,
                        "font_size_pt": float(run.font.size.pt) if run.font.size else None,
                        "color_hex": color_hex,
                    }
                )

            alignment = None
            if paragraph.alignment is not None:
                alignment = self._REVERSE_ALIGN_MAP.get(paragraph.alignment)

            paragraphs.append(
                {
                    "text": paragraph.text,
                    "level": paragraph.level,
                    "alignment": alignment,
                    "runs": runs,
                }
            )
        return paragraphs

    def _find_shape(self, slide, identifier):
        """Find a shape by name (str) or shape_id (int)."""
        # Try as integer shape_id first
        try:
            shape_id = int(identifier)
            for shape in slide.shapes:
                if int(shape.shape_id) == shape_id:
                    return shape
                # Also search inside groups
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP and hasattr(shape, "shapes"):
                    for child in shape.shapes:
                        if int(child.shape_id) == shape_id:
                            return child
        except (ValueError, TypeError):
            pass

        # Try as name
        name = str(identifier)
        for shape in slide.shapes:
            if shape.name == name:
                return shape
            # Also search inside groups
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP and hasattr(shape, "shapes"):
                for child in shape.shapes:
                    if child.name == name:
                        return child

        # Build available list for error
        available = [f"{s.name} (id={s.shape_id})" for s in slide.shapes]
        raise BridgeError(
            code="not_found",
            message=f"Shape '{identifier}' not found.",
            details={"identifier": str(identifier), "available_shapes": available},
        )

    def _detailed_shape_payload(self, shape) -> dict[str, Any]:
        """Return rich details about a shape for get_shape_details."""
        payload: dict[str, Any] = {
            "shape_id": int(shape.shape_id),
            "name": shape.name,
            "shape_type": str(shape.shape_type),
            "left_inches": emu_to_inches(shape.left) if shape.left is not None else None,
            "top_inches": emu_to_inches(shape.top) if shape.top is not None else None,
            "width_inches": emu_to_inches(shape.width) if shape.width is not None else None,
            "height_inches": emu_to_inches(shape.height) if shape.height is not None else None,
            "rotation": float(shape.rotation) if hasattr(shape, "rotation") else 0.0,
            "is_placeholder": bool(getattr(shape, "is_placeholder", False)),
        }

        # Text content
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            payload["has_text"] = True
            payload["text"] = shape.text
            payload["paragraphs"] = self._extract_paragraphs(shape.text_frame)
        else:
            payload["has_text"] = False

        # Table content
        if hasattr(shape, "has_table") and shape.has_table:
            table = shape.table
            payload["has_table"] = True
            payload["table_rows"] = len(table.rows)
            payload["table_cols"] = len(table.columns)
            cells = []
            for r_idx in range(len(table.rows)):
                row_cells = []
                for c_idx in range(len(table.columns)):
                    cell = table.cell(r_idx, c_idx)
                    row_cells.append(self._extract_cell_data(cell))
                cells.append(row_cells)
            payload["table_cells"] = cells
        else:
            payload["has_table"] = False

        # Picture
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            payload["is_picture"] = True
            try:
                payload["image_content_type"] = shape.image.content_type
            except Exception:
                payload["image_content_type"] = "unknown"
        else:
            payload["is_picture"] = False

        # Group
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP and hasattr(shape, "shapes"):
            payload["is_group"] = True
            payload["child_count"] = len(shape.shapes)
            payload["children"] = [self._shape_payload(child) for child in shape.shapes]
        else:
            payload["is_group"] = False

        return payload

    def _extract_cell_data(self, cell) -> dict[str, Any]:
        """Extract text and formatting from a table cell."""
        paragraphs = self._extract_paragraphs(cell.text_frame) if cell.text_frame else []
        return {
            "text": cell.text or "",
            "paragraphs": paragraphs,
        }

    def _write_cell(self, cell, data: dict[str, Any]) -> None:
        """Write text and optional formatting to a table cell."""
        text = str(data.get("text", ""))
        cell.text = text

        # Apply formatting to the first paragraph/run
        if cell.text_frame and cell.text_frame.paragraphs:
            paragraph = cell.text_frame.paragraphs[0]
            if paragraph.runs:
                self._apply_run_formatting(paragraph.runs[0], data)

            # Cell fill
            fill_hex = data.get("fill_hex")
            if fill_hex:
                try:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor.from_string(normalize_color(str(fill_hex)))
                except Exception as exc:
                    print(f"[pptx] Warning: could not set cell fill: {exc}", file=sys.stderr)

    def _replace_in_shapes(self, shapes, find_text: str, replace_text: str, case_sensitive: bool) -> int:
        replacements = 0
        search_pattern = re.compile(
            re.escape(find_text) if case_sensitive else re.escape(find_text), 0 if case_sensitive else re.IGNORECASE
        )

        for shape in shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if search_pattern.search(run.text):
                            old_text = run.text
                            run.text = search_pattern.sub(replace_text, run.text)
                            if run.text != old_text:
                                replacements += 1
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if search_pattern.search(run.text):
                                        old_text = run.text
                                        run.text = search_pattern.sub(replace_text, run.text)
                                        if run.text != old_text:
                                            replacements += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                replacements += self._replace_in_shapes(shape.shapes, find_text, replace_text, case_sensitive)
        return replacements
